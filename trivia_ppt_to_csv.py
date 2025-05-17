from pptx import Presentation
from pathlib import Path
import pandas as pd
import ftfy

# Load PowerPoint file
folder = Path('/Users/darthpatel/Library/CloudStorage'
              '/GoogleDrive-aotnashville@gmail.com/My Drive/Trivia_Nashville_AoT')

# COLLECT ALL THE PPTX FILES
ppt_files = folder.glob('*.pptx')

skip_phrases = ['Trivia','Question','bartender','donating',
                '$Googleplex','No Cheating allowed']

data = []

# GO THROUGH EACH PPTX FILE
for ppt_file in ppt_files:
    print(f"Processing {ppt_file.name}...")
    # Create Presentation Object
    presentation = Presentation(ppt_file)

    # Loop through slides
    for idx, slide in enumerate(presentation.slides, start=1):
        question_text = ""
        answer_text = ""

        # Extract question from text boxes
        question_parts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                text = shape.text.strip()
                if any(phrase.lower() in text.lower() for phrase in skip_phrases):
                    continue
                question_parts.append(text)
        question_text = ' '.join(question_parts).strip()

        # MAKE SURE TEXT ISN'T EMPTY
        if question_text:
            # Fix mojibake i.e., corrupted characters
            question_text = ftfy.fix_text(question_text)
            # Extract answer from tspeaker notes (if any)
            notes_slide = slide.notes_slide if slide.has_notes_slide else None
            if notes_slide and notes_slide.notes_text_frame:
                answer_text = notes_slide.notes_text_frame.text.strip()
                answer_text = ftfy.fix_text(answer_text)

            data.append({
                "file": ppt_file.stem,
                "Question": question_text,
                "Answer (from notes)": answer_text
            })

# Save to a CSV or view in DataFrame
df = pd.DataFrame(data)
df.to_csv("trivia_questions.csv", index=False, encoding='utf-8-sig')
print('PPT Converted to CSV')
# print(df)